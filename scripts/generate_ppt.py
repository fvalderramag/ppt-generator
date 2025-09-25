#!/usr/bin/env python3
import locale
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt  
from PIL import Image

# Configurar la localización en español
try:
    locale.setlocale(locale.LC_TIME, "es_ES.UTF-8")  # Linux/Mac
except:
    try:
        locale.setlocale(locale.LC_TIME, "Spanish_Spain")  # Windows
    except:
        print("⚠️ No se pudo establecer locale en español, se usará inglés.")


def parse_slides(md_path):
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    raw_slides = content.split("---")
    slides = []

    for raw in raw_slides:
        raw_lines = raw.strip().splitlines()
        if not raw_lines:
            continue

        title = None
        bullets = []
        images = []
        layout = "contenido"

        for line in raw_lines:
            stripped = line.strip()
            if not stripped:
                continue

            if stripped.startswith("#"):
                if "[layout:" in stripped:
                    parts = stripped.split("[layout:")
                    title = parts[0].lstrip("#").strip()
                    layout = parts[1].rstrip("]").strip()
                else:
                    title = stripped.lstrip("#").strip()
            elif stripped.endswith((".png", ".jpg", ".jpeg", ".drawio")):
                images.append(stripped)
            else:
                bullets.append(stripped.lstrip("-").strip())

        if title:
            slides.append((title, bullets, images, layout))
    return slides


def add_slide(prs, layout_map, title, bullets, images, layout_key):
    layout_index = layout_map.get(layout_key, layout_map["contenido"])
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])

    # ✅ Portada especial
    if layout_key == "portada":
        if slide.shapes.title:
            tf = slide.shapes.title.text_frame
            tf.clear()

            # Título
            p1 = tf.add_paragraph()
            p1.text = title

            # Espacio vacío
            tf.add_paragraph().text = ""

            # Fecha
            fecha = datetime.now().strftime("%B %Y").capitalize()
            tf.add_paragraph().text = fecha
        return

    # ✅ Layout imagen (solo título + imagen)
    if layout_key == "imagen":
        if slide.shapes.title:
            slide.shapes.title.text = title

        if images:
            for shape in slide.placeholders:
                if "Picture" in shape.name or "Imagen" in shape.name:                
                    ph = shape
                    # Tamaño del placeholder
                    ph_width, ph_height = ph.width, ph.height

                    # Tamaño real de la imagen
                    img = Image.open(images[0])
                    img_width, img_height = img.size

                    # Escala proporcional
                    ratio = min(ph_width / img_width, ph_height / img_height)
                    new_width = int(img_width * ratio)
                    new_height = int(img_height * ratio)

                    # Centrar la imagen dentro del placeholder
                    left = ph.left + (ph.width - new_width) // 2
                    top = ph.top + (ph.height - new_height) // 2

                    slide.shapes.add_picture(images[0], left, top, width=new_width, height=new_height)
                    break
        return

    # ✅ Otros layouts: título + bullets
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Buscar placeholder de cuerpo de texto (BODY = 2)
    body_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:
            body_shape = shape
            break

    if body_shape and bullets:
        tf = body_shape.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            run = p.add_run()
            run.text = bullet

            # 🔹 Reglas de formato
            if title.lower().startswith("agenda"):
                run.font.bold = True
                run.font.size = Pt(24)  # Tamaño Agenda
            elif title.lower().startswith("objetivos"):
                run.font.bold = False
                run.font.size = Pt(18)  # Tamaño Objetivos


def main():
    prs = Presentation("templates/template.pptx")
    slides = parse_slides("slides.md")

    layout_map = {
        "portada": 0,
        "contenido": 1,
        "imagen": 2,
    }

    for title, bullets, images, layout in slides:
        add_slide(prs, layout_map, title, bullets, images, layout)

    prs.save("presentacion.pptx")
    print("✅ Presentación generada con estilos y tamaños en Agenda/Objetivos: presentacion.pptx")


if __name__ == "__main__":
    main()
